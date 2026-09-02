#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成 GranularExtract 分析与说明 PPT。

依赖：python-pptx、matplotlib（安装在 managed venv）。
用法：python3 tools/gen_granular_ppt.py
输出：docs/GranularExtract_分析与说明.pptx
"""
import os
import math
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
DOC_DIR = os.path.join(ROOT, "docs")
ASSET_DIR = os.path.join(DOC_DIR, "_ppt_assets")
os.makedirs(ASSET_DIR, exist_ok=True)

# ---- 主题色 ----
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x2E, 0x6D, 0xA4)
TEAL = RGBColor(0x17, 0xA2, 0xB8)
AMBER = RGBColor(0xE8, 0x8A, 0x1A)
GREY = RGBColor(0x55, 0x5F, 0x6B)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x28, 0x2E)


def log2_ceil(x: int) -> int:
    return 0 if x <= 1 else math.ceil(math.log2(x))


def mux2_tree(N, M, G):
    K = log2_ceil((N - M) // G + 1)
    return K * M + G * ((1 << K) - 1 - K)


def mux2_bitmap(N, M, G):
    S = (N - M) // G + 1
    K = log2_ceil(S)
    return M * (S - 1) + S * K


# =========================================================================
# 图 1：多配置面积对比（B / T2 / T3≈N·K）
# =========================================================================
configs = [
    ("64/32/16", 64, 32, 16),
    ("64/32/8", 64, 32, 8),
    ("256/64/32", 256, 64, 32),
    ("512/128/64", 512, 128, 64),
    ("512/96/8", 512, 96, 8),
    ("1024/64/32", 1024, 64, 32),
    ("1024/256/64", 1024, 256, 64),
]
labels = [c[0] for c in configs]
b_vals = [mux2_tree(c[1], c[2], c[3]) for c in configs]
t2_vals = [mux2_bitmap(c[1], c[2], c[3]) for c in configs]
t3_vals = [c[1] * log2_ceil((c[1] - c[2]) // c[3] + 1) for c in configs]

fig, ax = plt.subplots(figsize=(11, 4.6))
import numpy as np
x = np.arange(len(labels))
w = 0.27
ax.bar(x - w, b_vals, w, label="B Windowed barrel tree", color="#2E6DA4")
ax.bar(x, t2_vals, w, label="T2 bitmap AND-OR plane", color="#E88A1A")
ax.bar(x + w, t3_vals, w, label="T3 full-width barrel (approx N*K)", color="#B0B7BE")
ax.set_yscale("log")
ax.set_ylabel("Area (equiv. mux2, log scale)")
ax.set_title("Area comparison across N/M/G configs (uniform mux2 unit)", fontsize=13, color="#1F3A5F", weight="bold")
ax.set_xticks(x)
ax.set_xticklabels(labels, rotation=20, ha="right")
ax.legend()
for i, v in enumerate(b_vals):
    ax.text(i - w, v * 1.1, str(v), ha="center", va="bottom", fontsize=8)
fig.tight_layout()
fig1 = os.path.join(ASSET_DIR, "chart_area.png")
fig.savefig(fig1, dpi=150)
plt.close(fig)

# =========================================================================
# 图 2：T2/B 比值 vs S（交叉点）
# =========================================================================
Ss = list(range(2, 70))
# 固定 N=512,M=96,G=8 → S=(512-96)/8+1=53；但为看趋势，用 M,G 固定、N 随 S 变
# 用 N=S*8+96-8 ... 简化：以 M=96,G=8,N 使 S 变化
ratios = []
for S in Ss:
    G = 8
    M = 96
    N = (S - 1) * G + M
    ratios.append(mux2_bitmap(N, M, G) / mux2_tree(N, M, G))

fig, ax = plt.subplots(figsize=(11, 4.4))
ax.plot(Ss, ratios, color="#17A2B8", linewidth=2.2, marker="o", ms=3)
ax.axhline(1.0, color="#E88A1A", linestyle="--", linewidth=1.4, label="ratio=1 (equal)")
ax.axvspan(0, 5, color="#E88A1A", alpha=0.10)
ax.set_xlabel("Number of selectable windows  S = (N-M)/G + 1")
ax.set_ylabel("T2 / B area ratio")
ax.set_title("bitmap plane vs barrel tree: T2 wins only at very small S", fontsize=13, color="#1F3A5F", weight="bold")
ax.annotate("T2 wins only here\n(S<=5 and large G)",
            xy=(4, 0.83), xytext=(14, 0.55),
            arrowprops=dict(arrowstyle="->", color="#E88A1A"),
            color="#E88A1A", fontsize=10)
ax.legend()
fig.tight_layout()
fig2 = os.path.join(ASSET_DIR, "chart_cross.png")
fig.savefig(fig2, dpi=150)
plt.close(fig)


# =========================================================================
# 构建 PPT
# =========================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, lines, size=18, color=DARK, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(ln, tuple):
            txt, kw = ln
        else:
            txt, kw = ln, {}
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(kw.get("size", size))
        run.font.bold = kw.get("bold", bold)
        run.font.color.rgb = kw.get("color", color)
        run.font.name = font
        if "space_after" in kw:
            p.space_after = Pt(kw["space_after"])
    return tb


def title_bar(slide, title, sub=None):
    rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), SW, Inches(0.06), TEAL)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), SW - Inches(1), Inches(0.9))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Microsoft YaHei"
    if sub:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(13); r2.font.color.rgb = RGBColor(0xCF, 0xDD, 0xEC); r2.font.name = "Microsoft YaHei"


def bullets(slide, x, y, w, h, items, size=16, gap=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple):
            head, body = it
        else:
            head, body = None, it
        if head:
            r = p.add_run(); r.text = "▸ " + head
            r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = BLUE; r.font.name = "Microsoft YaHei"
            r2 = p.add_run(); r2.text = body
            r2.font.size = Pt(size); r2.font.color.rgb = DARK; r2.font.name = "Microsoft YaHei"
        else:
            r = p.add_run(); r.text = "• " + body
            r.font.size = Pt(size); r.font.color.rgb = DARK; r.font.name = "Microsoft YaHei"
    return tb


# ---- 封面 ----
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.5), SW, Inches(0.08), TEAL)
textbox(s, Inches(0.8), Inches(2.1), SW - Inches(1.6), Inches(1.4),
        ["GranularExtract", ("按 G 粒度从 N bit 取 M bit 的自动选型模块", {"size": 20, "color": TEAL})],
        size=46, color=WHITE, bold=True)
textbox(s, Inches(0.8), Inches(3.8), SW - Inches(1.6), Inches(2.0),
        [("窗口化二分树 vs bitmap AND-OR 平面：按参数自动选择面积更小的实现", {"size": 18, "color": RGBColor(0xCF,0xDD,0xEC)}),
         ("配套：GranularExtract.scala · GranularExtractSpec（15 用例全绿）· 全量回归 266 用例通过", {"size": 14, "color": GREY, "space_after": 6}),
         ("BaseCbb / align  ·  2026-08-28", {"size": 14, "color": GREY})],
        size=18, color=WHITE)

# ---- 1. 问题定义 ----
s = add_slide()
title_bar(s, "1. 问题定义与约束", "从 N bit 输入按 G 粒度移位取出 M bit")
bullets(s, Inches(0.6), Inches(1.5), Inches(7.4), Inches(5.5), [
    ("起点对齐：", "start = s·G，s ∈ [0, S)，S = (N−M)/G + 1（可选窗口数）"),
    ("硬约束：", "N、M 均为 G 的整数倍 → N=n·G，M=m·G"),
    ("等价重述：", "从 n 个 G-bit chunk 中选 m 个连续 chunk"),
    ("输出语义：", "out[j] = in[start + j]，j ∈ [0, M)"),
    ("黄金模型：", "out = (in >> (off·G)) 的低 M bit"),
    ("典型场景：", "512b MAC 总线按 8B 切 96b 信元（N=512,M=96,G=8,S=53,K=6）"),
], size=17, gap=12)

# 右侧公式卡
rect(s, Inches(8.3), Inches(1.6), Inches(4.4), Inches(5.0), LIGHT)
textbox(s, Inches(8.6), Inches(1.8), Inches(3.9), Inches(4.7),
        [("参数关系", {"size": 18, "bold": True, "color": NAVY}),
         ("S = (N − M)/G + 1", {"size": 16, "color": DARK, "space_after": 6}),
         ("K = ⌈log₂ S⌉  （off 位宽）", {"size": 16, "color": DARK, "space_after": 6}),
         ("n = N / G，m = M / G", {"size": 16, "color": DARK, "space_after": 6}),
         ("约束：n ≥ m ≥ 1", {"size": 16, "color": DARK, "space_after": 6}),
         ("调用方须保证 off < S", {"size": 14, "color": AMBER})],
        size=16)

# ---- 2. 方案总览 ----
s = add_slide()
title_bar(s, "2. 候选方案一览", "统一用等效 2:1 mux 数（mux2）计量")
# 表格
rows = [
    ["方案", "结构", "面积 (mux2)", "组合深度", "适用"],
    ["B 二分块树", "K 级保持/下移 2^k·G", "K·M+G(2^K−1−K)", "K 级", "默认"],
    ["T2 bitmap 平面", "off 译码 S 路 one-hot", "M(S−1)+S·K", "译码+AND+OR", "S 极小且 G 大"],
    ["T3 整宽桶形", "整条 N 总线 K 级", "≈ N·K", "K 级", "不应单独用"],
    ["T1 字面AND-fold", "mask+AND+折叠", "≈ 2N", "log₂M", "仅滑窗匹配"],
    ["E/F 流式/SRAM", "串行交付/可寻址", "→ 0", "—", "本就流式"],
]
nrows, ncols = len(rows), 5
tbl_w = Inches(12.3)
gtbl = s.shapes.add_table(nrows, ncols, Inches(0.5), Inches(1.5), tbl_w, Inches(4.8)).table
gtbl.columns[0].width = Inches(2.2)
gtbl.columns[1].width = Inches(3.2)
gtbl.columns[2].width = Inches(2.9)
gtbl.columns[3].width = Inches(1.8)
gtbl.columns[4].width = Inches(2.2)
for ci in range(ncols):
    cell = gtbl.cell(0, ci)
    cell.text = rows[0][ci]
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]; p.runs[0].font.size = Pt(13); p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE; p.runs[0].font.name = "Microsoft YaHei"
for ri in range(1, nrows):
    for ci in range(ncols):
        cell = gtbl.cell(ri, ci)
        cell.text = rows[ri][ci]
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT if ri % 2 else WHITE
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(12); p.runs[0].font.color.rgb = DARK; p.runs[0].font.name = "Microsoft YaHei"
textbox(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7),
        [("T3 结论：循环移位 = 整宽桶形，恒比 B 贵约 N/M 倍（对整条总线移位，B 只对收缩窗口）。模块不实现 T3。",
          {"size": 13, "color": AMBER, "bold": True})], size=13)

# ---- 3. 面积对比图 ----
s = add_slide()
title_bar(s, "3. 面积对比（统一 mux2 单位）", "B 在大 S 下优势急剧放大")
s.shapes.add_picture(fig1, Inches(0.4), Inches(1.4), width=Inches(12.5))
textbox(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(1.0),
        [("B 与 T2 用同一 mux2 单位：512/96/8 下 T2 比 B 贵 5.15×；仅 64/32/16、512/128/64 等极小 S 大 G 场景 T2 略优。",
          {"size": 13, "color": GREY})], size=13)

# ---- 4. 交叉点 ----
s = add_slide()
title_bar(s, "4. 交叉点分析", "修正此前结论：T2 仅在极小 S 胜出")
s.shapes.add_picture(fig2, Inches(0.4), Inches(1.4), width=Inches(12.5))
textbox(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(1.1),
        [("此前未加约束的分析用「B×3、T2×1」的混合单位误判交叉点在 S≈7–10；统一单位后正确结论："
          "B 在绝大多数现实参数下更小，T2 仅在 S≤5 且 G 较大时反超。",
          {"size": 13, "color": GREY})], size=13)

# ---- 5. 多组对比表 ----
s = add_slide()
title_bar(s, "5. 多组 N/M/G 代价对比", "统一 mux2 单位 · 自动选型列")
tbl_rows = [
    ["N/M/G", "n,m,S(K)", "B", "T2", "T2/B", "自动选"],
    ["64/32/16", "4,2,3(2)", "80", "70", "0.875", "T2"],
    ["64/32/8", "8,4,5(3)", "128", "143", "1.117", "B"],
    ["128/32/8", "16,4,13(4)", "216", "436", "2.02", "B"],
    ["256/64/32", "8,2,7(3)", "320", "405", "1.27", "B"],
    ["512/128/64", "8,2,5(3)", "640", "527", "0.82", "T2"],
    ["512/96/8", "64,12,53(6)", "1032", "5310", "5.15", "B"],
    ["512/64/8", "64,8,57(6)", "840", "3926", "4.67", "B"],
    ["1024/64/32", "32,2,31(5)", "1152", "2075", "1.80", "B"],
    ["1024/256/64", "16,4,13(4)", "1728", "3124", "1.81", "B"],
]
nrows = len(tbl_rows)
gtbl = s.shapes.add_table(nrows, 6, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.3)).table
widths = [Inches(2.2), Inches(2.4), Inches(1.6), Inches(1.8), Inches(1.6), Inches(2.7)]
for ci, w in enumerate(widths):
    gtbl.columns[ci].width = w
for ci in range(6):
    cell = gtbl.cell(0, ci); cell.text = tbl_rows[0][ci]
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]; p.runs[0].font.size = Pt(13); p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE; p.runs[0].font.name = "Microsoft YaHei"
for ri in range(1, nrows):
    for ci in range(6):
        cell = gtbl.cell(ri, ci); cell.text = tbl_rows[ri][ci]
        cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if ri % 2 else WHITE
        p = cell.text_frame.paragraphs[0]; p.runs[0].font.size = Pt(12)
        p.runs[0].font.color.rgb = DARK if ci != 5 else (TEAL if tbl_rows[ri][5] == "T2" else BLUE)
        p.runs[0].font.bold = (ci == 5); p.runs[0].font.name = "Microsoft YaHei"

# ---- 6. 模块设计 ----
s = add_slide()
title_bar(s, "6. GranularExtractAuto 模块设计", "BaseCbb/align · 复用 GenModule")
bullets(s, Inches(0.6), Inches(1.5), Inches(7.2), Inches(5.5), [
    ("自动决策：", "elaboration 期估算 mux2Tree 与 mux2Bitmap，选更小者实例化，零运行时开销"),
    ("prefer 覆盖：", "auto / tree / bitmap，用于时序与布线调优"),
    ("IO：", "in(N) / off(K) / out(M) / sideIn(n) / sideOut(m)"),
    ("sideband：", "每 chunk 1 位标志（valid/last），宽 n→m，走同一套按 chunk 移位网络"),
    ("一致输出：", "B 与 T2 数学等价，均实现 out=(in>>off·G)(M-1,0)"),
    ("可流水：", "B 每级 2:1 mux，级间插 RegNext 把组合深度降到 1 级/拍（延迟+K，吞吐不变）"),
    ("非法 off：", "in 高位补 0、bitmap 越界项置 0，均为 don't-care 不耗面积"),
], size=15, gap=9)

rect(s, Inches(8.1), Inches(1.6), Inches(4.6), Inches(5.0), LIGHT)
textbox(s, Inches(8.35), Inches(1.8), Inches(4.1), Inches(4.7),
        [("面积估计（mux2）", {"size": 16, "bold": True, "color": NAVY}),
         ("mux2Tree =", {"size": 14, "color": DARK, "space_after": 2}),
         ("  K·M + G·(2^K−1−K)", {"size": 14, "color": BLUE, "bold": True, "space_after": 8}),
         ("mux2Bitmap =", {"size": 14, "color": DARK, "space_after": 2}),
         ("  M·(S−1) + S·K", {"size": 14, "color": AMBER, "bold": True, "space_after": 8}),
         ("useTree =", {"size": 14, "color": DARK, "space_after": 2}),
         ("  mux2Tree ≤ mux2Bitmap", {"size": 14, "color": TEAL, "bold": True, "space_after": 8}),
         ("面积诊断复用", {"size": 13, "color": GREY, "space_after": 2}),
         ("BaseCbb.Area", {"size": 13, "color": GREY}),
         (".ProcessConfiguration", {"size": 13, "color": GREY})],
        size=14)

# ---- 7. 验证 ----
s = add_slide()
title_bar(s, "7. 验证策略", "GranularExtractSpec · 15 用例全绿 · 全量回归 266 通过")
bullets(s, Inches(0.6), Inches(1.5), Inches(12.0), Inches(5.5), [
    ("① 黄金模型回环：", "8 组 N/M/G × 200 随机 off，断言 out == (in>>off·G)(M-1,0)"),
    ("② 自动选型断言：", "512/96/8→tree、64/32/8→tree、64/32/16→bitmap"),
    ("③ 两实现一致性：", "tree vs bitmap 随机比对，断言逐位相同（512/96/8、64/32/8）"),
    ("④ sideband 对齐：", "256/64/32 下 sideOut 按 chunk 移位与黄金模型一致"),
    ("⑤ prefer 覆盖：", "prefer=tree/bitmap 强制生效"),
    ("回归命令：", "sbt \"testOnly BaseCbb.align.GranularExtractSpec\""),
], size=16, gap=12)
rect(s, Inches(0.6), Inches(6.0), Inches(12.0), Inches(0.9), LIGHT)
textbox(s, Inches(0.8), Inches(6.05), Inches(11.6), Inches(0.8),
        [("全工程 sbt test：43 suites / 266 用例，0 失败（基线 251 + 新增 15）。",
          {"size": 14, "bold": True, "color": NAVY})], size=14)

# ---- 8. 落地建议 ----
s = add_slide()
title_bar(s, "8. 落地与选型建议", "数据通路中如何选")
bullets(s, Inches(0.6), Inches(1.6), Inches(12.0), Inches(5.2), [
    ("数据整块到达 + 偏移运行时变：", "用 GranularExtractAuto（默认 auto）"),
    ("S 极小（相邻 chunk、S≤3）且 G 较大：", "评估 prefer=bitmap（省 10–20%）"),
    ("输入是 G/cycle 流 / 已落 chunk 可寻址存储：", "选 E/F（选择逻辑趋零）"),
    ("检测而非提取（滑窗匹配）：", "选 T1（≈2N 有损 AND-fold），比「提取+判决」省"),
    ("不要加整宽循环移位步骤：", "T3 恒比 B 贵约 N/M 倍"),
    ("FPGA 备注：", "LUT6 原生 4:1 mux 缩小 B/T2 差距；SRL 让流式方案近乎免费"),
], size=16, gap=12)
rect(s, Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.7), NAVY)
textbox(s, Inches(0.8), Inches(6.55), Inches(11.6), Inches(0.6),
        [("模块已接入 BaseCbb/align，复用 GenModule 的 desiredName/参数化约定，可直接被寄存器生成工具与顶层整合。",
          {"size": 13, "color": WHITE})], size=13)

out = os.path.join(DOC_DIR, "GranularExtract_分析与说明.pptx")
prs.save(out)
print("SAVED:", out)
